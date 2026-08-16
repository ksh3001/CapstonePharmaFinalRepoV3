"""Build docs/product/aegis-one-pager.pptx — infographic one-pager."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "product" / "aegis-one-pager.pptx"
OUT_ALT = ROOT / "docs" / "product" / "aegis-one-pager-infographic.pptx"

PURPLE = RGBColor(0x3B, 0x1D, 0x6E)
PURPLE2 = RGBColor(0x6F, 0x42, 0xC1)
INK = RGBColor(0x2C, 0x33, 0x44)
MUTED = RGBColor(0x7B, 0x84, 0x94)
LINE = RGBColor(0xE6, 0xE8, 0xEE)
BG = RGBColor(0xF3, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x7A, 0x45)
BLUE = RGBColor(0x4D, 0xA6, 0xFF)
NAVY = RGBColor(0x3D, 0x5A, 0xFE)
GREEN = RGBColor(0x1F, 0x9D, 0x6A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LILAC = RGBColor(0xD4, 0xC4, 0xF0)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def rect(l: float, t: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None, rounded: bool = False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
        if line is not None:
            sh.line.fill.solid()
            sh.line.color.rgb = line
            sh.line.width = Pt(1)
        try:
            sh.shadow.inherit = False
        except Exception:
            pass
        return sh

    def text(l: float, t: float, w: float, h: float, value: str, *, size: int = 11, bold: bool = False, color: RGBColor = INK, align=PP_ALIGN.LEFT):
        sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = value
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return sh

    def lines(l: float, t: float, w: float, h: float, items: list[str], *, size: int = 11, color: RGBColor = INK, spacing: int = 2):
        sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = sh.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(spacing)
            run = p.add_run()
            run.text = item
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
        return sh

    def kpi(l: float, t: float, w: float, h: float, value: str, label: str, accent: RGBColor):
        rect(l, t, w, h, WHITE, LINE, rounded=True)
        rect(l, t, 0.12, h, accent)
        text(l + 0.22, t + 0.10, w - 0.35, 0.55, value, size=28, bold=True, color=accent)
        text(l + 0.22, t + 0.65, w - 0.35, 0.55, label, size=11, color=MUTED)

    rect(0, 0, 13.333, 7.5, BG)
    rect(0, 0, 13.333, 1.05, PURPLE)
    rect(0, 7.15, 13.333, 0.35, PURPLE)

    text(0.35, 0.14, 9.0, 0.42, "AEGIS", size=30, bold=True, color=WHITE)
    text(
        0.35,
        0.58,
        10.2,
        0.35,
        "Advisory HITL evidence console  ·  Engines classify  ·  Humans decide",
        size=13,
        color=RGBColor(0xE8, 0xDC, 0xFF),
    )
    text(10.0, 0.28, 3.0, 0.55, "Batch  ·  PV  ·  Supply", size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    for x, val, lab, acc in (
        (0.30, "84", "Injects covered\n81 rule + 3 artefact", GREEN),
        (2.45, "358", "Automated tests\nunittest suite", PURPLE2),
        (4.60, "7", "Runtime roles\nentitlement matrix", ORANGE),
        (6.75, "6", "Declared agents\nAG-1 … AG-6", BLUE),
        (8.90, "14", "Feature specs\n146 BR · 204 AC", NAVY),
        (11.05, "3", "Workflows live\n11 catalog ids", TEAL),
    ):
        kpi(x, 1.20, 2.05, 1.25, val, lab, acc)

    text(0.35, 2.60, 8.0, 0.28, "Three mandated workflows", size=14, bold=True, color=PURPLE)
    for x, acc, title, ids, body in (
        (0.30, PURPLE2, "Batch evidence", "NCB204-B24071", "Quality holds · contradictions\nForced evidence before ack"),
        (4.55, ORANGE, "PV intake", "PV-1001 · SM-77", "Clocks · duplicates · listedness\nNever confirms a signal"),
        (8.80, BLUE, "Supply / cold-chain", "SH-901 · shortage", "Draft options only\nNo stock movement"),
    ):
        rect(x, 2.95, 4.05, 1.35, WHITE, LINE, rounded=True)
        rect(x, 2.95, 4.05, 0.10, acc)
        text(x + 0.18, 3.12, 3.7, 0.28, title, size=14, bold=True, color=acc)
        text(x + 0.18, 3.42, 3.7, 0.24, ids, size=11, bold=True, color=INK)
        lines(x + 0.18, 3.68, 3.7, 0.55, body.split("\n"), size=11, color=MUTED, spacing=1)

    rect(0.30, 4.45, 4.25, 2.55, WHITE, LINE, rounded=True)
    text(0.48, 4.55, 3.9, 0.28, "7 runtime roles", size=14, bold=True, color=PURPLE)
    lines(
        0.48,
        4.90,
        3.9,
        1.95,
        [
            "•  EU Qualified Person",
            "•  Safety physician",
            "•  Supply governance",
            "•  Quality reviewer",
            "•  CISO / DPO",
            "•  Auditor (read-only)",
            "•  Unblinding authority",
        ],
        size=12,
        spacing=3,
    )

    rect(4.70, 4.45, 4.25, 2.55, WHITE, LINE, rounded=True)
    text(4.88, 4.55, 3.9, 0.28, "6 agents · control stack", size=14, bold=True, color=PURPLE)
    lines(
        4.88,
        4.90,
        3.9,
        1.95,
        [
            "AG-1 Supervisor  ·  approve interrupt",
            "AG-2 Evidence retrieval",
            "AG-3 Batch  ·  AG-4 PV  ·  AG-5 Supply",
            "AG-6 Review packager",
            "",
            "Deny-list  ·  kill switch  ·  residency",
            "23 control-map rows  ·  11 ADRs",
            "5 modes: assessment → cloud",
        ],
        size=12,
        spacing=2,
    )

    rect(9.10, 4.45, 3.90, 2.55, WHITE, LINE, rounded=True)
    text(9.28, 4.55, 3.55, 0.28, "Run in 3 steps", size=14, bold=True, color=PURPLE)
    lines(
        9.28,
        4.90,
        3.55,
        1.35,
        [
            "1  pip install -r requirements-ui.txt",
            "2  python -m aegis setup",
            "3  python -m aegis serve --port 8000",
            "",
            "http://127.0.0.1:8000/",
        ],
        size=11,
        spacing=2,
    )
    text(
        9.28,
        6.35,
        3.55,
        0.50,
        "Not a disposition / PV-decision /\nstock-movement / recall system",
        size=10,
        bold=True,
        color=ORANGE,
    )

    text(
        0.35,
        7.20,
        12.6,
        0.25,
        "AEGIS  ·  FDE capstone one-pager  ·  84/84 injects  ·  358 tests  ·  synthetic fixtures  ·  advisory posture (not certified)  ·  2026",
        size=10,
        color=LILAC,
    )

    try:
        prs.save(OUT)
        return OUT
    except PermissionError:
        prs.save(OUT_ALT)
        return OUT_ALT


if __name__ == "__main__":
    path = build()
    print(path)
