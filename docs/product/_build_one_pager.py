"""Build the AEGIS one-page infographic. Not imported by the app."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("aegis-one-pager.pptx")

PURPLE = RGBColor(0x3B, 0x1D, 0x6E)
PURPLE2 = RGBColor(0x6F, 0x42, 0xC1)
INK = RGBColor(0x2C, 0x33, 0x44)
MUTED = RGBColor(0x7B, 0x84, 0x94)
LINE = RGBColor(0xE6, 0xE8, 0xEE)
BG = RGBColor(0xF3, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x7A, 0x45)
BLUE = RGBColor(0x4D, 0xA6, 0xFF)
NAVY = RGBColor(0x3D, 0x5A, 0xFE)


def _shape(slide, kind, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def _set_text(shape, lines, *, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    for i, (text, *opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts[2] if len(opts) > 2 else align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts[0] if opts else size)
        run.font.bold = opts[1] if len(opts) > 1 else bold
        run.font.color.rgb = opts[3] if len(opts) > 3 else color
        run.font.name = "Calibri"


def kpi(slide, l, t, w, h, value, label, accent):
    card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, WHITE, LINE)
    card.adjustments[0] = 0.12
    bar = _shape(slide, MSO_SHAPE.RECTANGLE, l, t, 0.09, h, accent)
    _set_text(
        card,
        [(value, 28, True, accent, PP_ALIGN.CENTER), (label, 11, False, INK, PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
    )


def chip(slide, l, t, w, h, title, sub, body, accent):
    card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, WHITE, LINE)
    card.adjustments[0] = 0.08
    _shape(slide, MSO_SHAPE.RECTANGLE, l, t, 0.10, h, accent)
    _set_text(
        card,
        [
            (title, 14, True, PURPLE, PP_ALIGN.LEFT),
            (sub, 10, True, accent, PP_ALIGN.LEFT),
            (body, 11, False, INK, PP_ALIGN.LEFT),
        ],
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, BG)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.08, PURPLE)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.22, 13.333, 0.28, PURPLE)

    title = _shape(slide, MSO_SHAPE.RECTANGLE, 0.32, 0.10, 8.4, 0.50, PURPLE)
    _set_text(title, [("AEGIS", 30, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    sub = _shape(slide, MSO_SHAPE.RECTANGLE, 0.32, 0.58, 9.4, 0.40, PURPLE)
    _set_text(
        sub,
        [("Advisory HITL evidence console  ·  engines classify  ·  humans decide", 13, False, RGBColor(0xE8, 0xDC, 0xFF))],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    badge = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.15, 0.28, 2.85, 0.52, PURPLE2)
    _set_text(badge, [("FDE capstone  ·  2026", 12, True, WHITE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    kpis = [
        (0.32, "84", "Injects covered\nINJ-001 … 084", PURPLE2),
        (2.90, "410", "Automated tests\n189 test modules", ORANGE),
        (5.48, "7", "Runtime roles\nQP, safety, supply…", BLUE),
        (8.06, "14", "Feature specs\nFR-001 … FR-014", NAVY),
        (10.64, "3", "Workflows\nBatch · PV · Supply", PURPLE),
    ]
    for x, value, label, accent in kpis:
        kpi(slide, x, 1.24, 2.40, 1.28, value, label.replace("\n", "  ·  "), accent)

    mid = [
        (0.32, "146", "Business rules"),
        (2.46, "204", "Acceptance criteria"),
        (4.60, "27", "NFRs"),
        (6.74, "6", "Declared agents"),
        (8.88, "11", "ADRs"),
        (11.02, "23", "Compliance controls"),
    ]
    for x, value, label in mid:
        card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.66, 1.98, 0.78, WHITE, LINE)
        card.adjustments[0] = 0.16
        _set_text(
            card,
            [(value, 20, True, PURPLE, PP_ALIGN.CENTER), (label, 10, False, MUTED, PP_ALIGN.CENTER)],
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
        )

    chip(
        slide,
        0.32,
        3.60,
        4.20,
        1.55,
        "Batch evidence",
        "NCB204-B24071",
        "Quality-hold pack. Forced evidence view. Acknowledge is a workflow event, not a signature.",
        PURPLE2,
    )
    chip(
        slide,
        4.62,
        3.60,
        4.20,
        1.55,
        "PV intake",
        "PV-1001  ·  SM-77",
        "Clocks and duplicate candidates retained. Entitlement blinds sensitive segments. No signal confirmation.",
        ORANGE,
    )
    chip(
        slide,
        8.92,
        3.60,
        4.10,
        1.55,
        "Supply / cold-chain",
        "SH-901  ·  shortage",
        "Draft options and holds only. Contested capacity is not available stock.",
        BLUE,
    )

    left = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.32, 5.30, 6.35, 1.78, WHITE, LINE)
    _set_text(
        left,
        [
            ("Console surface", 13, True, PURPLE),
            ("Dashboard · contradictions · status · evidence history", 12, False, INK),
            ("User guide (header ?) with live-screen pictures", 12, False, INK),
            ("Pack chat: catalog status only — decide-language refused", 12, False, INK),
            ("Oversight: gates · inject map · 6 agents · runtime health", 12, False, INK),
        ],
    )
    right = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.82, 5.30, 6.20, 1.78, WHITE, LINE)
    _set_text(
        right,
        [
            ("Run  ·  CPython 3.11–3.13  ·  no Node  ·  no DB", 13, True, PURPLE),
            ("pip install -r requirements-ui.txt", 12, False, INK),
            ("python -m aegis setup   then   serve --port 8000", 12, False, INK),
            ("http://127.0.0.1:8000/     assessment = offline, model off", 12, False, INK),
            ("C4: docs/architecture/aegis-architecture.html", 12, False, INK),
        ],
    )

    foot = _shape(slide, MSO_SHAPE.RECTANGLE, 0.32, 7.22, 12.7, 0.28, PURPLE)
    _set_text(
        foot,
        [
            (
                "Not a disposition / PV-decision / eligibility / stock-movement / recall system   ·   "
                "Not EU AI Act or ISO 42001 certified   ·   7 runtime roles, AI authority none",
                10,
                False,
                RGBColor(0xD4, 0xC4, 0xF0),
                PP_ALIGN.LEFT,
            )
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
